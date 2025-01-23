import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from typing import List, Dict
from dataclasses import dataclass

# Assuming your LayoutModel is defined in a module named 'models'
# and your layout algorithm in a module named 'layout'
# Update these imports according to your project structure
from bcm.models import LayoutModel
from bcm.layout_manager import process_layout
from bcm.settings import Settings


@dataclass
class NodeSize:
    width: float
    height: float


@dataclass
class GridLayout:
    rows: int
    cols: int
    width: float
    height: float
    deviation: float
    positions: List[Dict[str, float]]


def calculate_node_size(node: LayoutModel, settings: Settings) -> NodeSize:
    """Calculate the minimum size needed for a node and its children."""
    if not node.children:
        return NodeSize(settings.get("box_min_width"), settings.get("box_min_height"))

    child_sizes = [calculate_node_size(child, settings) for child in node.children]
    best_layout = find_best_layout(child_sizes, len(node.children), settings)

    return NodeSize(best_layout.width, best_layout.height)


def find_best_layout(
    child_sizes: List[NodeSize], child_count: int, settings: Settings
) -> GridLayout:
    """
    Find the optimal grid layout for a set of child nodes.
    """
    best_layout = GridLayout(
        rows=1,
        cols=child_count,
        width=float("inf"),
        height=float("inf"),
        deviation=float("inf"),
        positions=[],
    )

    horizontal_gap = settings.get("horizontal_gap")
    vertical_gap = settings.get("vertical_gap")
    padding = settings.get("padding")
    top_padding = settings.get("top_padding", padding)  # Correctly get top_padding
    target_aspect_ratio = settings.get("target_aspect_ratio")

    for rows_tentative in range(1, child_count + 1):
        for cols_float in [
            child_count / rows_tentative,
            (child_count + rows_tentative - 1) // rows_tentative,
        ]:
            cols = int(round(cols_float))
            if cols == 0:
                continue
            rows = (child_count + cols - 1) // cols

            row_heights = [0.0] * rows
            col_widths = [0.0] * cols

            # Calculate maximum heights and widths for each row and column
            for i in range(child_count):
                row = i // cols
                col = i % cols
                size = child_sizes[i]

                row_heights[row] = max(row_heights[row], size.height)
                col_widths[col] = max(col_widths[col], size.width)

            grid_width = sum(col_widths) + (cols - 1) * horizontal_gap
            grid_height = sum(row_heights) + (rows - 1) * vertical_gap

            # Calculate total dimensions including padding
            total_width = grid_width + 2 * padding
            total_height = (
                grid_height + top_padding + padding
            )  # Use top_padding and bottom padding

            aspect_ratio = total_width / total_height
            deviation = abs(aspect_ratio - target_aspect_ratio)

            # Calculate positions for each child
            positions = []
            y_offset = top_padding  # Start at top_padding

            # Calculate extra space for distributing among rows and columns
            extra_width_per_col = (
                max(0, total_width - (grid_width + 2 * padding)) / cols
                if cols > 0
                else 0
            )
            extra_height_per_row = (
                max(0, total_height - (grid_height + top_padding + padding)) / rows
                if rows > 0
                else 0
            )

            for row in range(rows):
                x_offset = padding
                for col in range(cols):
                    idx = row * cols + col
                    if idx < child_count:
                        child_position = {
                            "x": x_offset,
                            "y": y_offset,
                            "width": col_widths[col] + extra_width_per_col,
                            "height": row_heights[row] + extra_height_per_row,
                        }
                        positions.append(child_position)
                        x_offset += (
                            col_widths[col] + extra_width_per_col + horizontal_gap
                        )
                y_offset += row_heights[row] + extra_height_per_row + vertical_gap

            # **Calculate the actual height needed based on the last child's position**
            max_child_bottom = 0
            for pos in positions:
                max_child_bottom = max(max_child_bottom, pos["y"] + pos["height"])
            actual_height = max_child_bottom + padding

            current_layout = GridLayout(
                rows=rows,
                cols=cols,
                width=total_width,
                height=actual_height,  # Use the actual height
                deviation=deviation,
                positions=positions,
            )

            if current_layout.deviation < best_layout.deviation or (
                current_layout.deviation == best_layout.deviation
                and current_layout.width * current_layout.height
                < best_layout.width * best_layout.height
            ):
                best_layout = current_layout

    return best_layout


def layout_tree(
    node: LayoutModel, settings: Settings, x: float = 0, y: float = 0
) -> LayoutModel:
    """Recursively layout the tree starting from the given node."""
    if not node.children:
        node.width = settings.get("box_min_width")
        node.height = settings.get("box_min_height")
        node.x = x
        node.y = y
        return node

    layout = find_best_layout(
        [calculate_node_size(child, settings) for child in node.children],
        len(node.children),
        settings,
    )

    node.width = layout.width
    node.height = layout.height
    node.x = x
    node.y = y

    for child, pos in zip(node.children, layout.positions):
        layout_tree(child, settings, x + pos["x"], y + pos["y"])
        child.width = pos["width"]
        child.height = pos["height"]

    return node


def hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


def pixels_to_inches(pixels: float, scale_factor: float = 1.0) -> float:
    """Convert pixels to inches (assuming 96 DPI) and apply scaling."""
    return (pixels / 96.0) * scale_factor


def calculate_font_size(
    root_size: int, level: int, is_leaf: bool, scale_factor: float = 1.0
) -> int:
    """Calculate font size based on level, node type, and apply scaling."""
    base_size = root_size - (level * 4)
    if is_leaf:
        font_size = max(base_size - 2, 8)
    else:
        font_size = max(base_size - 4, 10)
    # Scale down the font size and ensure it's not less than 1
    scaled_font_size = max(int(font_size * scale_factor), 1)
    return scaled_font_size


def add_node_to_group(
    parent_group,
    node: LayoutModel,
    settings: Settings,
    scale_factor: float,
    level: int = 0,
):
    """Add a node and its children to the group shape."""
    # Convert coordinates and dimensions to inches, applying scaling
    left = pixels_to_inches(node.x, scale_factor)
    top = pixels_to_inches(node.y, scale_factor)
    width = pixels_to_inches(node.width, scale_factor)
    height = pixels_to_inches(node.height, scale_factor)

    # Determine node color based on level and whether it has children
    if not node.children:
        color = settings.get("color_leaf")
    else:
        color = settings.get(f"color_{min(level, 6)}")

    # Convert hex color to RGB
    rgb_color = hex_to_rgb(color)

    # Add group shape for the current node and its children
    group = parent_group.shapes.add_group_shape()

    # Add shape to group
    shape = group.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,  # Rectangle shape
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )

    # Calculate the adjustment value
    smaller_dimension = min(width, height)
    adjustment_value = 0.015 / smaller_dimension

    # To ensure that the radius does not exceed the maximum possible curvature, we limit the radius to 0.5 or 50% of the smallest dimension
    if adjustment_value > 0.5:
        adjustment_value = 0.5
        print("Radius exceeds maximum allowed. Setting to maximum.")

    shape.adjustments[0] = adjustment_value

    # Set shape fill color
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb_color)

    # Set shape line color
    shape.line.color.rgb = RGBColor(51, 51, 51)  # #333333
    shape.line.width = Pt(0.2)

    # Calculate font size with scaling
    font_size = calculate_font_size(
        settings.get("root_font_size"), level, not node.children, scale_factor
    )

    # Add text
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = None
    text_frame.vertical_anchor = MSO_ANCHOR.TOP if node.children else MSO_ANCHOR.MIDDLE

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.text = node.name

    # Set font properties
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0, 0, 0)

    # Recursively add child nodes
    if node.children:
        for child in node.children:
            add_node_to_group(group, child, settings, scale_factor, level + 1)


def export_to_pptx(
    model: LayoutModel, settings: Settings, scale_factor: float = 0.3
) -> Presentation:
    """
    Export the capability model to PowerPoint format.

    :param model: The root LayoutModel of the capability model.
    :param settings: A Settings object containing configuration.
    :param scale_factor: A factor to scale down the layout.
    :return: A Presentation object.
    """
    # Process layout
    processed_model = process_layout(model, settings)

    # Create presentation
    prs = Presentation()

    # Add slide
    slide_layout = prs.slide_layouts[6]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)

    # Calculate dimensions with padding and apply scaling
    padding = settings.get("padding", 20)
    scaled_padding = padding * scale_factor
    width = math.ceil(processed_model.width + 2 * padding) * scale_factor
    height = math.ceil(processed_model.height + 2 * padding) * scale_factor

    # Set slide dimensions based on the scaled width and height
    prs.slide_width = Inches(pixels_to_inches(width, 1))
    prs.slide_height = Inches(pixels_to_inches(height, 1))

    # Add group shape to slide
    group_shape = slide.shapes.add_group_shape()

    # Add nodes recursively with scaling
    add_node_to_group(group_shape, processed_model, settings, scale_factor)

    return prs
