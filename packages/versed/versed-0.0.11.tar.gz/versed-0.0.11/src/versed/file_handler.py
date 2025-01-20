import csv
import io
import json
from pathlib import Path
from typing import Dict

from docx import Document
from docx.text.paragraph import Paragraph
from docx.table import Table
import pandas as pd
from pptx import Presentation
from pypdf import PdfReader
from versed.gdrive_file_handler import GoogleDriveHandler

class LocalFileHandler:

    def __init__(self):
        pass

    def _get_local_file_stream(self, file: Dict) -> Dict:
        """
        Gets the file stream of a file from disk.

        Args:
            file (dict): {"name": file_name, "path": file path on disk }

        Returns:
            Dict: {
                "stream": io.BytesIO: The file stream of the file,
                "type": str: The file extension of the file
            }
        """
        extension = Path(file["name"]).suffix
        return {
            "stream": self._get_file_stream(file),
            "type": extension,
        }

    def _get_file_stream(self, file: Dict) -> io.BytesIO:
        """
        Gets the file stream of a file from disk.

        Args:
            file (dict): {"path": file path on disk}

        Returns:
            io.BytesIO: The file stream of the file.
        """
        file_path = Path(file["path"])
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        with open(file_path, "rb") as f:
            return io.BytesIO(f.read())


class FileHandler:

    def __init__(self, credentials):
        self.local_file_handler = LocalFileHandler()
        self.gdrive_file_handler = GoogleDriveHandler(credentials)

    def get_file_stream(self, file) -> io.BytesIO:
        if file["path"].startswith("gdrive://"):
            resolved_file = self.gdrive_file_handler._get_google_drive_file_stream(file)
        else:
            resolved_file = self.local_file_handler._get_local_file_stream(file)

        return resolved_file
    
    def get_file_content(self, file) -> str:
        resolved_file = self.get_file_stream(file)
        extension = resolved_file["type"]

        match extension:
            case ".txt":
                content = self._get_txt_content(resolved_file)
            case ".docx" | ".gdoc":
                content = self._get_docx_content(resolved_file)
            case ".pdf":
                content = self._get_pdf_content(resolved_file)
            case ".csv" | ".gsheet":
                content = self._get_csv_content(resolved_file)
            case ".pptx" | ".gslides":
                content = self._get_pptx_content(resolved_file)
            case ".ipynb":
                content = self._get_notebook_content(resolved_file)
            case ".xlsx":
                content = self._get_excel_content(resolved_file)
            case _:
                raise ValueError(f"Unsupported file type: {extension}")

        return content
    
    def _get_txt_content(self, file: Dict) -> str:
        """
        Extracts string content from a .txt file.
        """
        return file["stream"].read().decode("utf-8")

    def _get_docx_content(self, file: Dict) -> str:
        """
        Extracts string content from a Word (.docx) file.
        """
        doc = Document(file["stream"])

        file_content = ""

        # Iterate over all block-level elements
        for element in doc.iter_inner_content():
            if isinstance(element, Paragraph):
                file_content += f"{element.text}\n\n"
            elif isinstance(element, Table):
                file_content += "<Table>\n"
                for row in element.rows:
                    file_content += "<Row>\n"
                    for cell in row.cells:
                        file_content += "<Cell>\n"
                        file_content += f"{cell.text}\n"
                        file_content += "</Cell>\n"
                    file_content += "</Row>\n"
                file_content += "</Table>\n\n"

        return file_content

    def _get_pdf_content(self, file: Dict) -> str:
        """
        Extracts string content from a PDF file.
        """
        reader = PdfReader(file["stream"])
        content = ""
        for page in reader.pages:
            content += page.extract_text() + "\n"
        return content

    def _get_csv_content(self, file: Dict) -> str:
        """
        Extracts string content from a CSV file.
        """
        content = ""
        reader = csv.reader(io.TextIOWrapper(file["stream"], encoding="utf-8"))
        for row in reader:
            content += ",".join(row) + "\n"
        return content

    def _get_pptx_content(self, file: Dict) -> str:
        """
        Extracts string content from a PowerPoint (.pptx) file.
        """
        presentation = Presentation(file["stream"])
        content = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
        return content

    def _get_notebook_content(self, file: Dict) -> str:
        """
        Extracts string content from a Jupyter Notebook (.ipynb) file, including all cell types.

        Args:
            file (Dict): The notebook file dictionary.

        Returns:
            str: The extracted content from the notebook.
        """
        # Load the notebook JSON content
        notebook = json.load(file["stream"])
        content = ""

        # Iterate over all cells in the notebook
        for cell in notebook.get("cells", []):
            cell_type = cell.get("cell_type")
            cell_source = "".join(cell.get("source", []))
            
            if cell_type == "code":
                content += f"[Code Cell]:\n{cell_source}\n\n"
            elif cell_type == "markdown":
                content += f"[Markdown Cell]:\n{cell_source}\n\n"
            elif cell_type == "raw":
                content += f"[Raw Cell]:\n{cell_source}\n\n"
            else:
                content += f"[Unknown Cell Type]:\n{cell_source}\n\n"

        return content

    def _get_excel_content(self, file: Dict) -> str:
        """
        Extracts string content from an Excel (.xlsx) file using pandas.

        Args:
            file (Dict): A file dictionary.

        Returns:
            str: The extracted content from all sheets in the Excel file.
        """
        excel_data = pd.ExcelFile(file["stream"])
        content = ""

        # Iterate over all sheets
        for sheet_name in excel_data.sheet_names:
            content += f"Sheet: {sheet_name}\n"
            sheet_df = excel_data.parse(sheet_name)
            content += sheet_df.to_string(index=False)  # Convert DataFrame to a readable string
            content += "\n\n"

        return content
