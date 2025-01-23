from __future__ import annotations

import abc
import asyncio
import io
import logging
import os
import subprocess
import tempfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Never, Optional, Sequence, TypedDict, cast, override
from xml.etree.ElementTree import Element

import msgspec
from architecture import dp, log
from architecture.data.files import FileExtension, RawFile, find_extension
from architecture.utils.decorators import ensure_module_installed
from architecture.utils.functions import run_sync
from openai import OpenAI

from intellibricks.agents import Agent
from intellibricks.llms.types import (
    AudioDescription,
    AudioFilePart,
    ChainOfThought,
    ImageFilePart,
    MimeType,
    VideoFilePart,
    VisualMediaDescription,
)

from .constants import ParsingStrategy
from .parsed_files import Image, ParsedFile, SectionContent, TablePageItem

debug_logger = log.create_logger(__name__, level=logging.DEBUG)
exception_logger = log.create_logger(__name__, level=logging.ERROR)


class InvalidFileExtension(Exception):
    """Raised when a file extension is not supported."""


class LocalSettings(TypedDict):
    use_gpu: bool


class FileParser(msgspec.Struct, frozen=True, tag_field="type"):
    """
    Abstract class for extracting content from files.
    This should be used as a base class for specific file parsers.
    """

    strategy: ParsingStrategy = ParsingStrategy.DEFAULT

    def extract_contents(
        self,
        file: RawFile,
    ) -> ParsedFile:
        """Extracts content from the file."""
        return run_sync(self.extract_contents_async, file)

    @abc.abstractmethod
    async def extract_contents_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        """Extracts content from the file."""
        raise NotImplementedError("This method should be implemented by subclasses.")


@dp.Facade
class IntellibricksFileParser(FileParser, frozen=True, tag="intellibricks"):
    visual_description_agent: Optional[
        Agent[ChainOfThought[VisualMediaDescription]]
    ] = None
    """Agent used for generating textual descriptions of images and videos, if the synapse supports it."""

    audio_description_agent: Optional[Agent[ChainOfThought[AudioDescription]]] = None
    """Agent used for generating textual descriptions of audio files, if the synapse supports it."""

    @override
    async def extract_contents_async(self, file: RawFile) -> ParsedFile:
        match file.extension:
            # Word files
            case FileExtension.DOC | FileExtension.DOCX:
                debug_logger.debug("Extracting contents from Word file")
                return await OfficeFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)

            # PowerPoint
            case FileExtension.PPT | FileExtension.PPTX | FileExtension.PPTM:
                debug_logger.debug("Extracting contents from PowerPoint file")
                return await OfficeFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)

            # Excel
            case FileExtension.XLS | FileExtension.XLSX:
                debug_logger.debug("Extracting contents from Excel file")
                return await OfficeFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)

            case FileExtension.TXT:
                debug_logger.debug("Extracting contents from TXT file")
                return await TxtFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)

            # Treat XML as plain text for now
            case FileExtension.XML:
                debug_logger.debug("Extracting contents from XML file")
                return await XMLFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)

            case FileExtension.PDF:
                debug_logger.debug("Extracting contents from PDF file")
                return await PDFFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)

            # Static images (PNG, JPG, TIFF, BMP)
            case (
                FileExtension.JPEG
                | FileExtension.PNG
                | FileExtension.TIFF
                | FileExtension.BMP
                | FileExtension.JPG
            ):
                debug_logger.debug("Extracting contents from static image file")
                return await StaticImageFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)

            case FileExtension.GIF:
                debug_logger.debug("Extracting contents from animated image file")
                return await AnimatedImageFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)

            case FileExtension.PKT:
                debug_logger.debug("Extracting contents from PKT file")
                return await PKTFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)

            case FileExtension.ALG:
                debug_logger.debug("Extracting contents from ALG file")
                return await AlgFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)

            case FileExtension.ZIP | FileExtension.RAR | FileExtension.PKZ:
                debug_logger.debug("Extracting contents from compressed file")
                return await CompressedFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)

            case FileExtension.DWG:
                debug_logger.debug("Extracting contents from DWG file")
                return await DWGFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)
            case (
                FileExtension.FLAC
                | FileExtension.MP3
                | FileExtension.MPEG
                | FileExtension.MPGA
                | FileExtension.M4A
                | FileExtension.OGG
                | FileExtension.WAV
                | FileExtension.WEBM
            ):
                debug_logger.debug("Extracting contents from audio file")
                return await AudioFileParser(
                    strategy=self.strategy,
                    audio_description_agent=self.audio_description_agent,
                ).extract_contents_async(file)
            case FileExtension.MP4:
                debug_logger.debug("Extracting contents from video file")
                return await VideoFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)
            case _:
                raise InvalidFileExtension(
                    f"Unsupported file extension: {file.extension}"
                )


class XMLFileParser(IntellibricksFileParser, frozen=True, tag="xml"):
    """
    Parses XML files, keeping the raw XML in the 'text' field and converting the structure
    into a nested Markdown list in the 'md' field.
    """

    @override
    async def extract_contents_async(self, file: RawFile) -> ParsedFile:
        raw_xml = file.contents.decode("utf-8", errors="replace")
        md_content = self.xml_to_md(raw_xml)

        section_content = SectionContent(
            number=1,
            text=raw_xml,
            md=md_content,
            images=[],
            items=[],
        )

        return ParsedFile(
            name=file.name,
            sections=[section_content],
        )

    def xml_to_md(self, xml_str: str) -> str:
        """Converts XML content into a nested Markdown list structure."""
        try:
            root: Element = ET.fromstring(xml_str)
            return self._convert_element_to_md(root, level=0)
        except ET.ParseError as e:
            exception_logger.exception("Error parsing XML: %s", e)
            return "```xml\n" + xml_str + "\n```"  # Fallback to raw XML in code block

    def _convert_element_to_md(self, element: Element, level: int) -> str:
        """Recursively converts an XML element and its children to Markdown.

        Args:
            element: The XML element to convert
            level: Current nesting level for indentation
        """
        indent = "  " * level
        lines: list[str] = []

        # Element tag as bold item
        lines.append(f"{indent}- **{element.tag}**")

        # Attributes as sub-items
        if element.attrib:
            lines.append(f"{indent}  - *Attributes*:")
            for key, value in element.attrib.items():
                lines.append(f"{indent}    - `{key}`: `{value}`")

        # Text content
        if element.text and element.text.strip():
            text = element.text.strip().replace("\n", " ")
            lines.append(f"{indent}  - *Text*: {text}")

        # Process child elements recursively
        for child in element:
            lines.append(self._convert_element_to_md(child, level + 1))

        return "\n".join(lines)


class CompressedFileParser(IntellibricksFileParser, frozen=True, tag="compressed"):
    """
    Parses compressed files (ZIP, RAR, PKZ) by extracting each file within the archive,
    delegating to the appropriate parser, and merging the results.
    """

    @override
    async def extract_contents_async(self, file: RawFile) -> ParsedFile:
        import tempfile
        import zipfile

        import rarfile

        # We'll accumulate ParsedFile objects from each extracted child file
        parsed_files: list[ParsedFile] = []

        # Write the compressed file to a temporary location
        with tempfile.NamedTemporaryFile(delete=True) as tmp:
            tmp.write(file.contents)
            tmp.flush()

            # Decide how to open the archive based on extension
            match file.extension:
                case FileExtension.ZIP | FileExtension.PKZ:
                    # Treat PKZ exactly like ZIP for demo purposes
                    with zipfile.ZipFile(tmp.name, "r") as zip_ref:
                        # Iterate over files inside the archive
                        for info in zip_ref.infolist():
                            # Directories have filename ending with "/"
                            if info.is_dir():
                                continue

                            # Read raw bytes of the child file
                            child_data = zip_ref.read(info)
                            child_name = info.filename
                            child_ext = find_extension(filename=child_name)

                            # Turn that child file into a RawFile
                            child_raw_file = RawFile.from_bytes(
                                data=child_data,
                                name=child_name,
                                extension=child_ext,
                            )

                            # Parse using our IntellibricksFileParser façade
                            # (re-using the same strategy/visual_description_agent)
                            parser = IntellibricksFileParser(
                                strategy=self.strategy,
                                visual_description_agent=self.visual_description_agent,
                            )
                            child_parsed = await parser.extract_contents_async(
                                child_raw_file
                            )
                            parsed_files.append(child_parsed)

                case FileExtension.RAR:
                    with rarfile.RarFile(tmp.name, "r") as rar_ref:
                        for info in rar_ref.infolist():
                            """Type of "isdir" is unknownPylancereportUnknownMemberType"""
                            if info.isdir():  # type: ignore
                                continue

                            child_data = rar_ref.read(info)  # type: ignore

                            child_name = info.filename  # type: ignore

                            child_ext = find_extension(filename=child_name)  # type: ignore

                            child_raw_file = RawFile.from_bytes(
                                data=child_data,  # type: ignore
                                name=child_name,  # type: ignore
                                extension=child_ext,
                            )

                            parser = IntellibricksFileParser(
                                strategy=self.strategy,
                                visual_description_agent=self.visual_description_agent,
                            )
                            child_parsed = await parser.extract_contents_async(
                                child_raw_file
                            )
                            parsed_files.append(child_parsed)

                case _:
                    # Fallback if something else accidentally calls this parser
                    raise ValueError(
                        f"CompressedFileParser does not handle extension: {file.extension}"
                    )

        # Merge all the parsed files into a single ParsedFile
        return ParsedFile.from_parsed_files(parsed_files)


class DWGFileParser(IntellibricksFileParser, frozen=True, tag="dwg"):
    @ensure_module_installed("aspose-cad", "intellibricks[files]")
    @override
    async def extract_contents_async(self, file: RawFile) -> ParsedFile:
        """
        DWG files are kind of tricky. To parse them, Intellibricks converts them to PDF first,
        then takes a "screenshot" of each page of the PDF and uses GenAI to describe the images.
        """
        import platform

        if platform.machine() == "arm64":
            raise ValueError("ARM architecture is not supported by aspose-cad")

        import aspose.cad as cad  # type: ignore

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            # Load the DWG file
            image = cad.Image.load(file_path)  # type: ignore

            # Specify PDF Options
            pdfOptions = cad.imageoptions.PdfOptions()  # type: ignore

            output_path = f"{temp_dir}/output.pdf"

            # Save as PDF
            image.save(output_path, pdfOptions)  # type: ignore

            image_bytes_list = self.__pdf_to_images(output_path)

            raw_files = [
                RawFile.from_bytes(
                    data=img, name=f"{file.name}_{i}.png", extension=FileExtension.PNG
                )
                for i, img in enumerate(image_bytes_list)
            ]

            parser = StaticImageFileParser(
                strategy=self.strategy,
                visual_description_agent=self.visual_description_agent,
            )

            parsed_files = [await parser.extract_contents_async(f) for f in raw_files]
            sections = [
                section
                for parsed_file in parsed_files
                for section in parsed_file.sections
            ]

            return ParsedFile.from_sections(file.name, sections)

    def __pdf_to_images(self, pdf_path: str) -> Sequence[bytes]:
        """Converts each page of a PDF to image bytes.

        Args:
            pdf_path (str): The path to the PDF file.

        Returns:
            Sequence[bytes]: A list of bytes objects, each containing a PNG image of a PDF page.
        """
        import pymupdf

        image_bytes_list: list[bytes] = []
        doc = pymupdf.open(pdf_path)

        try:
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)  # type: ignore
                pix = page.get_pixmap()  # type: ignore

                # Create a bytes buffer and save the image into it
                buffer = io.BytesIO()
                pix.save(buffer, "png")  # type: ignore
                image_bytes = buffer.getvalue()

                image_bytes_list.append(image_bytes)

        finally:
            doc.close()

        return image_bytes_list


class PKTFileParser(IntellibricksFileParser, frozen=True, tag="pkt"):
    @override
    async def extract_contents_async(self, file: RawFile) -> ParsedFile:
        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            xml_bytes = self.pkt_to_xml_bytes(file_path)

            # For now, we'll just return the XML content as a single page
            xml_text = xml_bytes.decode("utf-8", errors="replace")

            page_content = SectionContent(
                number=1,
                text=xml_text,
                md=xml_text,
                images=[],
                items=[],
            )

            return ParsedFile(
                name=file.name,
                sections=[page_content],
            )

    def pkt_to_xml_bytes(self, pkt_file: str) -> bytes:
        """
        Convert a Packet Tracer file (.pkt/.pka) to its XML representation as bytes.

        :param pkt_file: Path to the input .pkt or .pka file.
        :return: The uncompressed XML content as bytes.
        """
        import zlib

        with open(pkt_file, "rb") as f:
            in_data = bytearray(f.read())

        i_size = len(in_data)
        out = bytearray()

        # Decrypt each byte with decreasing file length
        for byte in in_data:
            out.append(byte ^ (i_size & 0xFF))
            i_size -= 1

        # The first 4 bytes (big-endian) represent the size of the XML when uncompressed
        # (This value is not needed for the actual return, but we parse it for completeness.)
        _uncompressed_size = int.from_bytes(out[:4], byteorder="big")

        # Decompress the data after the first 4 bytes
        xml_data = zlib.decompress(out[4:])

        return xml_data


class AlgFileParser(IntellibricksFileParser, frozen=True, tag="alg"):
    """ALG Files can be treated as text files, so we'll use TxtFileParser to extract content."""

    @override
    async def extract_contents_async(self, file: RawFile) -> ParsedFile:
        return await TxtFileParser(
            strategy=self.strategy,
            visual_description_agent=self.visual_description_agent,
        ).extract_contents_async(file)


class PDFFileParser(IntellibricksFileParser, frozen=True, tag="pdf"):
    @ensure_module_installed("pypdf", "intellibricks[files]")
    @override
    async def extract_contents_async(self, file: RawFile) -> ParsedFile:
        from pypdf import PdfReader

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            reader = PdfReader(file_path)
            section_contents: list[SectionContent] = []
            for page_num, page in enumerate(reader.pages):
                page_images: Sequence[Image] = [
                    Image(contents=image.data, name=image.name) for image in page.images
                ]

                image_descriptions: list[str] = []
                if (
                    self.visual_description_agent
                    and self.strategy == ParsingStrategy.HIGH
                ):
                    for image_num, image in enumerate(page_images):
                        agent_input = ImageFilePart(
                            mime_type=MimeType.image_png, data=image.contents
                        )
                        agent_response = await self.visual_description_agent.run_async(
                            agent_input
                        )
                        image_md: str = agent_response.parsed.final_answer.md
                        image_descriptions.append(
                            f"Page Image {image_num + 1}: {image_md}"
                        )

                page_text = [page.extract_text(), "".join(image_descriptions)]

                md = "".join(page_text)
                section_content = SectionContent(
                    number=page_num + 1,
                    text=md,
                    md=md,
                    images=page_images,
                )

                section_contents.append(section_content)

            file_name = file.name
            return ParsedFile(
                name=file_name,
                sections=section_contents,
            )


class OfficeFileParser(IntellibricksFileParser, frozen=True, tag="office"):
    """
    This class actually delegates the parsing to the appropriate parser based on the file extension.
    This class is a Facade for the different Office file parsers.
    """

    async def extract_contents_async(self, file: RawFile) -> ParsedFile:
        extension = file.extension
        match extension:
            # Word
            case FileExtension.DOC | FileExtension.DOCX:
                return await DocxFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)

            # PowerPoint (including .ppt, .pptx, .pptm)
            case FileExtension.PPT | FileExtension.PPTX | FileExtension.PPTM:
                return await PptxFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)

            # Excel (including .xls, .xlsx)
            case FileExtension.XLS | FileExtension.XLSX:
                return await ExcelFileParser(
                    strategy=self.strategy,
                    visual_description_agent=self.visual_description_agent,
                ).extract_contents_async(file)

            case _:
                raise ValueError(f"Unsupported Office extension: {extension}")


class DocxFileParser(OfficeFileParser, frozen=True, tag="docx"):
    @ensure_module_installed("docx", "intellibricks[files]")
    @override
    async def extract_contents_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        import tempfile

        from docx import Document  # python-docx

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            document = Document(file_path)

            # Extract all text from paragraphs
            paragraph_texts: list[str] = []
            for paragraph in document.paragraphs:
                if paragraph.text.strip():
                    paragraph_texts.append(paragraph.text)
            doc_text = "\n".join(paragraph_texts)

            # Extract all images
            doc_images: list[Image] = []
            for rel in document.part._rels.values():  # type: ignore
                # Relationship is image-based if it references an image part
                if "image" in rel.reltype:
                    image_part = rel.target_part
                    image_name = image_part.partname.split("/")[-1]  # e.g. "image1.png"
                    image_bytes = image_part.blob
                    doc_images.append(Image(name=image_name, contents=image_bytes))

            # If high-level strategy, describe images
            image_descriptions: list[str] = []
            if self.visual_description_agent and self.strategy == ParsingStrategy.HIGH:
                for idx, image in enumerate(doc_images, start=1):
                    agent_input = ImageFilePart(
                        mime_type=MimeType.image_png,  # or detect from extension
                        data=image.contents,
                    )
                    agent_response = await self.visual_description_agent.run_async(
                        agent_input
                    )
                    image_md = agent_response.parsed.final_answer.md
                    image_descriptions.append(f"Docx Image {idx}: {image_md}")

                # Append the images' descriptions to the main text
                if image_descriptions:
                    doc_text += "\n\n" + "\n".join(image_descriptions)

            # Create a single SectionContent (DOCX has no true "pages" by default)
            section_content = SectionContent(
                number=1,
                text=doc_text,
                md=doc_text,
                images=doc_images,
            )

            return ParsedFile(
                name=file.name,
                sections=[section_content],
            )


class PptxFileParser(OfficeFileParser, frozen=True, tag="pptx"):
    @ensure_module_installed("pptx", "intellibricks[files]")
    @override
    async def extract_contents_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.presentation import Presentation as PptxPresentation
        from pptx.shapes.autoshape import Shape
        from pptx.shapes.picture import Picture

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            if file.extension in {FileExtension.PPT, FileExtension.PPTM}:
                converted_pptx_file = self._convert_to_pptx(file)
                converted_pptx_file.save_to_file(file_path)
            else:
                file.save_to_file(file_path)

            prs: PptxPresentation = Presentation(file_path)

            sections: list[SectionContent] = []

            for slide_index, slide in enumerate(prs.slides, start=1):
                # We'll store text from shapes and images
                slide_texts: list[str] = []
                slide_images: list[Image] = []

                # Examine each shape
                for shape in slide.shapes:
                    # If shape has a text frame, cast to Shape
                    if shape.has_text_frame:
                        shape_with_text = cast(Shape, shape)
                        text_str: str = shape_with_text.text
                        slide_texts.append(text_str)

                    # If shape is a picture, cast to Picture
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        picture_shape = cast(Picture, shape)
                        image_blob: bytes = picture_shape.image.blob
                        image_name: str = shape.name or f"slide_{slide_index}_img"
                        slide_images.append(Image(name=image_name, contents=image_blob))

                combined_text: str = "\n".join(slide_texts)

                # If strategy is HIGH, we generate image descriptions
                if (
                    self.visual_description_agent
                    and self.strategy == ParsingStrategy.HIGH
                ):
                    image_descriptions: list[str] = []
                    for img_idx, image_obj in enumerate(slide_images, start=1):
                        agent_input = ImageFilePart(
                            mime_type=MimeType.image_png,
                            data=image_obj.contents,
                        )
                        agent_response = await self.visual_description_agent.run_async(
                            agent_input
                        )
                        image_md: str = agent_response.parsed.final_answer.md
                        image_descriptions.append(
                            f"Slide {slide_index} - Image {img_idx}: {image_md}"
                        )

                    if image_descriptions:
                        combined_text += "\n\n" + "\n".join(image_descriptions)

                section_content = SectionContent(
                    number=slide_index,
                    text=combined_text,
                    md=combined_text,
                    images=slide_images,
                )
                sections.append(section_content)

            return ParsedFile(
                name=file.name,
                sections=sections,
            )

    def _convert_to_pptx(self, file: RawFile) -> RawFile:
        """Convert PowerPoint files (.ppt/.pptm) to .pptx format and return as RawFile.

        Args:
            file: RawFile instance containing the input file data.

        Returns:
            RawFile instance containing converted content.

        Raises:
            RuntimeError: If conversion fails or LibreOffice not installed.
        """

        def _is_libreoffice_installed() -> bool:
            try:
                subprocess.run(
                    ["libreoffice", "--version"],
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                    check=True,
                )
                return True
            except (subprocess.CalledProcessError, FileNotFoundError):
                return False

        if not _is_libreoffice_installed():
            raise RuntimeError("LibreOffice not found in system PATH")

        with tempfile.TemporaryDirectory() as temp_dir:
            # Write input file to temporary directory
            input_path = os.path.join(temp_dir, file.name)
            with open(input_path, "wb") as f:
                f.write(file.contents)

            # Run LibreOffice conversion
            try:
                subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "pptx",
                        "--outdir",
                        temp_dir,
                        input_path,
                    ],
                    check=True,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.PIPE,
                    timeout=60,
                )
            except subprocess.CalledProcessError as e:
                error_msg = e.stderr.decode().strip() if e.stderr else "Unknown error"
                raise RuntimeError(f"Conversion failed: {error_msg}") from e
            except subprocess.TimeoutExpired:
                raise RuntimeError("Conversion timed out after 60 seconds")

            # Determine output file path
            output_filename = Path(file.name).stem + ".pptx"
            output_path = os.path.join(temp_dir, output_filename)

            if not os.path.exists(output_path):
                available_files = os.listdir(temp_dir)
                raise RuntimeError(
                    f"Converted file not found at {output_path}. Found files: {available_files}"
                )

            # Read converted file and return as RawFile
            return RawFile.from_file_path(output_path)


class ExcelFileParser(OfficeFileParser, frozen=True, tag="excel"):
    @ensure_module_installed("openpyxl", "intellibricks[files]")
    @override
    async def extract_contents_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        import csv
        import io

        from openpyxl import Workbook, load_workbook

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            wb: Workbook = load_workbook(file_path, data_only=True)
            sections: list[SectionContent] = []

            for sheet_index, sheet in enumerate(wb.worksheets, start=1):
                # Gather structured data
                rows: list[list[str]] = []
                row_texts: list[str] = []
                for row in sheet.iter_rows(values_only=True):
                    # Process cell values
                    cell_values = [
                        str(cell) if cell is not None else "" for cell in row
                    ]
                    rows.append(cell_values)
                    row_texts.append("\t".join(cell_values))

                combined_text = "\n".join(row_texts)

                # Generate CSV content
                csv_buffer = io.StringIO()
                csv_writer = csv.writer(csv_buffer)
                csv_writer.writerows(rows)
                csv_str = csv_buffer.getvalue().strip()

                # Process images
                sheet_images: list[Image] = []
                if hasattr(sheet, "_images"):
                    image_list = getattr(sheet, "_images", [])
                    for img_idx, img in enumerate(image_list, start=1):
                        img_data = getattr(img, "_data", None)
                        if img_data is not None:
                            image_name = f"{sheet.title}_img_{img_idx}.png"
                            sheet_images.append(
                                Image(name=image_name, contents=img_data)
                            )

                # Generate image descriptions if needed
                if (
                    self.visual_description_agent
                    and self.strategy == ParsingStrategy.HIGH
                ):
                    image_descriptions: list[str] = []
                    for img_idx, image_obj in enumerate(sheet_images, start=1):
                        agent_input = ImageFilePart(
                            mime_type=MimeType.image_png,
                            data=image_obj.contents,
                        )
                        agent_response = await self.visual_description_agent.run_async(
                            agent_input
                        )
                        image_md = agent_response.parsed.final_answer.md
                        image_descriptions.append(
                            f"Worksheet {sheet.title} - Image {img_idx}: {image_md}"
                        )
                    if image_descriptions:
                        combined_text += "\n\n" + "\n".join(image_descriptions)

                # Create table page item
                table_item = TablePageItem(
                    md=combined_text, rows=rows, csv=csv_str, is_perfect_table=True
                )

                section_content = SectionContent(
                    number=sheet_index,
                    text=combined_text,
                    md=combined_text,
                    images=sheet_images,
                    items=[table_item],
                )
                sections.append(section_content)

            return ParsedFile(
                name=file.name,
                sections=sections,
            )


class TxtFileParser(IntellibricksFileParser, frozen=True, tag="txt"):
    """
    Parses plain .txt files. Extracts all content as a single page (number=1).
    """

    @override
    async def extract_contents_async(self, file: RawFile) -> ParsedFile:
        text_content = file.contents.decode("utf-8", errors="replace")

        page_content = SectionContent(
            number=1,
            text=text_content,
            md=text_content,
        )

        return ParsedFile(
            name=file.name,
            sections=[page_content],
        )


class StaticImageFileParser(IntellibricksFileParser, frozen=True, tag="static_image"):
    """
    Parses static image files (PNG, JPEG, TIFF, etc.) as a single "page" with one image.
    If the image is TIFF, it converts to PNG in-memory for better compatibility.
    If the strategy == HIGH and an visual_description_agent is present,
    it appends an AI-generated textual description of the image.
    """

    @override
    async def extract_contents_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        from PIL import Image as PILImage

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            # Determine the extension
            extension = file.extension.value.lower()  # e.g. "png", "jpg", "tiff"

            # Convert to PNG if TIFF
            if extension in {"tiff", "tif"}:
                # Use Pillow to open, then convert to PNG in memory
                with io.BytesIO(file.contents) as input_buffer:
                    with PILImage.open(input_buffer) as pil_img:
                        # Convert to RGBA or RGB if needed
                        if pil_img.mode not in ("RGB", "RGBA"):
                            pil_img = pil_img.convert("RGBA")

                        # Save as PNG into a new buffer
                        output_buffer = io.BytesIO()
                        pil_img.save(output_buffer, format="PNG")
                        converted_bytes = output_buffer.getvalue()

                # Use the converted PNG bytes
                image_bytes = converted_bytes
                current_mime_type = MimeType.image_png
            else:
                # No conversion needed
                image_bytes = file.contents

                # For demonstration, pick your MIME by extension
                if extension in {"png", "bmp"}:
                    current_mime_type = MimeType.image_png
                elif extension in {"jpg", "jpeg"}:
                    current_mime_type = MimeType.image_jpeg
                else:
                    # Fallback to PNG or raise an error if you want
                    current_mime_type = MimeType.image_png

            # Create an Image object
            image_obj = Image(name=file.name, contents=image_bytes)

            # Generate a description if we have an agent + HIGH strategy
            text_content = ""
            if self.visual_description_agent and self.strategy == ParsingStrategy.HIGH:
                agent_input = ImageFilePart(
                    mime_type=current_mime_type,
                    data=image_bytes,
                )
                agent_response = await self.visual_description_agent.run_async(
                    agent_input
                )
                description_md = agent_response.parsed.final_answer.md
                text_content = description_md

            # We treat it as a single "page" with one image
            page_content = SectionContent(
                number=1,
                text=text_content,
                md=text_content,
                images=[image_obj],
            )

            return ParsedFile(
                name=file.name,
                sections=[page_content],
            )


class AnimatedImageFileParser(
    IntellibricksFileParser, frozen=True, tag="animated_image"
):
    """
    Parses animated GIF files by splitting them into 3 equally sized segments
    (or fewer if total frames < 3).
    Each selected frame is turned into a PNG in memory and, if strategy == HIGH,
    sent to visual_description_agent for a textual description.

    Returns a ParsedFile with up to 3 SectionContent items, each page representing
    one of the frames chosen from the animation.
    """

    @override
    async def extract_contents_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        from PIL import Image as PILImage

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            # Safety check: only proceed if it's a .gif
            # or you can attempt detection based on file headers
            extension = file.extension.value.lower()
            if extension not in {"gif"}:
                raise ValueError("AnimatedImageFileParser only supports .gif files.")

            # --- 1. Load all frames from the GIF ---
            frames: list[PILImage.Image] = []
            with PILImage.open(file_path) as gif_img:
                try:
                    while True:
                        frames.append(gif_img.copy())
                        gif_img.seek(gif_img.tell() + 1)
                except EOFError:
                    pass  # we've reached the end of the animation

            num_frames = len(frames)
            if num_frames == 0:
                # No frames => no content
                return ParsedFile(name=file.name, sections=[])

            # --- 2. Pick up to 3 frames, splitting the GIF into 3 segments ---
            # If there are fewer than 3 frames, just use them all.
            # If more than 3, pick three frames spaced across the animation.

            if num_frames <= 3:
                selected_frames = frames
            else:
                # Example approach: pick near 1/3, 2/3, end
                idx1 = max(0, (num_frames // 3) - 1)
                idx2 = max(0, (2 * num_frames // 3) - 1)
                idx3 = num_frames - 1
                # Ensure distinct indexes
                unique_indexes = sorted(set([idx1, idx2, idx3]))
                selected_frames = [frames[i] for i in unique_indexes]

            # --- 3. Convert each selected frame to PNG and (optionally) describe it ---
            pages: list[SectionContent] = []
            for i, frame in enumerate(selected_frames, start=1):
                # Convert frame to PNG in-memory
                png_buffer = io.BytesIO()
                # Convert to RGBA if needed
                if frame.mode not in ("RGB", "RGBA"):
                    frame = frame.convert("RGBA")
                frame.save(png_buffer, format="PNG")
                png_bytes = png_buffer.getvalue()

                # Create an Image object
                frame_image = Image(
                    name=f"{file.name}-frame{i}.png", contents=png_bytes
                )

                # If strategy is HIGH, pass the frame to the agent
                text_description = ""
                if (
                    self.visual_description_agent
                    and self.strategy == ParsingStrategy.HIGH
                ):
                    agent_input = ImageFilePart(
                        mime_type=MimeType.image_png,
                        data=png_bytes,
                    )
                    agent_response = await self.visual_description_agent.run_async(
                        agent_input
                    )
                    text_description = agent_response.parsed.final_answer.md

                # Each frame is its own "page" in the final doc
                page_content = SectionContent(
                    number=i,
                    text=text_description,
                    md=text_description,
                    images=[frame_image],
                )
                pages.append(page_content)

            # --- 4. Return the multi-page ParsedFile ---
            return ParsedFile(
                name=file.name,
                sections=pages,
            )


class AudioFileParser(IntellibricksFileParser, frozen=True, tag="audio"):
    async def extract_contents_async(self, file: RawFile) -> ParsedFile:
        if self.audio_description_agent is None:
            raise ValueError("No audio description agent provided.")

        file_contents: bytes = file.contents
        file_extension: FileExtension = file.extension

        if file_extension in {
            FileExtension.FLAC,
            FileExtension.MPEG,
            FileExtension.MPGA,
            FileExtension.M4A,
            FileExtension.OGG,
            FileExtension.WAV,
            FileExtension.WEBM,
        }:
            import aiofiles.os as aios
            from aiofiles import open as aio_open

            self._check_ffmpeg_installed()

            # Generate unique temporary filenames
            input_temp = os.path.join(
                tempfile.gettempdir(),
                f"input_{os.urandom(8).hex()}.{file_extension.value}",
            )
            output_temp = os.path.join(
                tempfile.gettempdir(), f"output_{os.urandom(8).hex()}.mp3"
            )

            # Write input file asynchronously
            async with aio_open(input_temp, "wb") as f:
                await f.write(file_contents)

            # Build FFmpeg command
            command = [
                "ffmpeg",
                "-hide_banner",
                "-loglevel",
                "error",  # Suppress unnecessary logs
                "-y",  # Overwrite output file if exists
                "-i",
                input_temp,
                "-codec:a",
                "libmp3lame",
                "-q:a",
                "2",  # Quality preset (0-9, 0=best)
                output_temp,
            ]

            # Execute FFmpeg
            process = await asyncio.create_subprocess_exec(
                *command, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
            )

            _, stderr = await process.communicate()

            # Handle conversion errors
            if process.returncode != 0:
                await aios.remove(input_temp)
                if await aios.path.exists(output_temp):
                    await aios.remove(output_temp)
                raise RuntimeError(
                    f"Audio conversion failed: {stderr.decode().strip()}"
                )

            # Read converted file
            async with aio_open(output_temp, "rb") as f:
                file_contents = await f.read()

            # Cleanup temporary files
            await aios.remove(input_temp)
            await aios.remove(output_temp)

        transcription = self.audio_description_agent.run(
            AudioFilePart(data=file_contents, mime_type=MimeType.audio_mp3)
        )

        return ParsedFile(
            name=file.name,
            sections=[
                SectionContent(
                    number=1,
                    text=transcription.raw_transcription
                    or self._could_not_transcript(),
                    md=transcription.parsed.final_answer.md,
                    images=[],
                )
            ],
        )

    def _could_not_transcript(self) -> Never:
        raise ValueError("Could not transcribe the audio")

    def _check_ffmpeg_installed(self) -> None:
        import subprocess

        try:
            result = subprocess.run(
                ["ffmpeg", "-version"],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
            )
            exception_logger.exception("FFmpeg is not installed or not in PATH.")
            if result.returncode != 0:
                raise RuntimeError()
        except FileNotFoundError:
            exception_logger.exception("FFmpeg is not installed or not in PATH.")
            raise RuntimeError()


class VideoFileParser(IntellibricksFileParser, frozen=True, tag="video"):
    async def extract_contents_async(self, file: RawFile) -> ParsedFile:
        if self.visual_description_agent is None:
            raise ValueError("No visual description agent provided.")

        extension = file.extension
        if extension != FileExtension.MP4:
            raise ValueError("VideoFileParser only supports .mp4 files.")

        file_contents = file.contents
        visual_media_description = await self.visual_description_agent.run_async(
            VideoFilePart(data=file_contents, mime_type=MimeType.video_mp4)
        )

        return ParsedFile(
            name=file.name,
            sections=[
                SectionContent(
                    number=1,
                    text=visual_media_description.parsed.final_answer.md,
                    md=visual_media_description.parsed.final_answer.md,
                    images=[],
                )
            ],
        )


class MarkitdownFileParser(FileParser, frozen=True, tag="markitdown"):
    client: Optional[OpenAI] = None
    model: Optional[str] = None

    @ensure_module_installed("markitdown", "intellibricks[files]")
    async def extract_contents_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        from markitdown import MarkItDown
        from markitdown._markitdown import DocumentConverterResult

        match self.strategy:
            case (
                ParsingStrategy.DEFAULT | ParsingStrategy.MEDIUM | ParsingStrategy.FAST
            ):
                llm_client = None
                llm_model = None
            case ParsingStrategy.HIGH:
                llm_client = self.client or OpenAI()
                llm_model = self.model or "gpt-4o"

        with tempfile.NamedTemporaryFile(delete=True) as temp_file:
            temp_file.write(file.contents)
            temp_file.seek(0)
            converter = MarkItDown(llm_client=llm_client, llm_model=llm_model)
            result: DocumentConverterResult = converter.convert(temp_file.name)
            markdown: str = result.text_content

            # return a Document with one page only
            page_content = SectionContent(
                number=1,
                text=markdown,
                md=markdown,
                images=[],
                items=[],
            )

            return ParsedFile(
                name=file.name,
                sections=[page_content],
            )
